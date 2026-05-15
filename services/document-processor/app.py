#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
RAG System - Document Processor Service（Ollama専用）
日本語対応PDFとOfficeファイルの処理サービス
"""

import asyncio
import hashlib
import json
import os
from contextlib import asynccontextmanager
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

# OCR ライブラリ
import easyocr
# HTTP クライアント
import httpx
import openpyxl
# ドキュメント処理用ライブラリ
import pypdf
import pytesseract
import uvicorn
from docx import Document
from fastapi import BackgroundTasks, FastAPI, File, HTTPException, UploadFile
from fastapi.responses import JSONResponse
# 日本語テキスト処理
from janome.tokenizer import Tokenizer
from loguru import logger
from pdf2image import convert_from_path
from PIL import Image
from pptx import Presentation
from sqlalchemy import create_engine, text
from sqlalchemy.orm import sessionmaker

# 設定
DATABASE_URL = os.getenv("DATABASE_URL")
USE_OLLAMA = os.getenv("USE_OLLAMA", "true").lower() == "true"

# OCRリーダーを初期化（日本語対応）
ocr_reader = easyocr.Reader(['ja', 'en'], gpu=False)
OLLAMA_URL = os.getenv("OLLAMA_URL", "http://ollama:11434")
EMBEDDING_MODEL = os.getenv("EMBEDDING_MODEL", "nomic-embed-text:latest")
UPLOAD_DIR = Path("/app/uploads")
PROCESSED_DIR = Path("/app/processed")

# psycopg3対応のデータベース接続
if DATABASE_URL:
    if "postgresql://" in DATABASE_URL:
        DATABASE_URL = DATABASE_URL.replace("postgresql://", "postgresql+psycopg://")

engine = create_engine(DATABASE_URL)
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)

# 日本語トークナイザー
tokenizer = Tokenizer()

def ensure_database_schema():
    """データベーススキーマを確認・更新"""
    try:
        with engine.connect() as conn:
            # filenameカラムの長さを確認・修正
            try:
                conn.execute(text("""
                    ALTER TABLE documents 
                    ALTER COLUMN filename TYPE VARCHAR(255)
                """))
                conn.commit()
                logger.info("データベーススキーマを更新しました（filename VARCHAR(255)）")
            except Exception as e:
                if "already exists" in str(e) or "does not exist" in str(e):
                    logger.info("データベーススキーマは既に最新です")
                else:
                    logger.warning(f"スキーマ更新警告: {e}")
            
            # file_typeカラムの長さも確認・修正
            try:
                conn.execute(text("""
                    ALTER TABLE documents 
                    ALTER COLUMN file_type TYPE VARCHAR(100)
                """))
                conn.commit()
                logger.info("データベーススキーマを更新しました（file_type VARCHAR(100)）")
            except Exception as e:
                if "already exists" in str(e) or "does not exist" in str(e):
                    logger.info("file_typeカラムは既に適切です")
                else:
                    logger.warning(f"file_typeスキーマ更新警告: {e}")
                    
    except Exception as e:
        logger.error(f"データベーススキーマ確認エラー: {e}")

# アプリ起動時にスキーマ確認


@asynccontextmanager
async def lifespan(app: FastAPI):
    """アプリケーションのライフサイクル管理"""
    # 起動時の処理
    logger.info("Document Processor Service 起動中...")
    ensure_database_schema()
    UPLOAD_DIR.mkdir(exist_ok=True)
    PROCESSED_DIR.mkdir(exist_ok=True)
    logger.info("Document Processor Service 起動完了")
    yield
    # 終了時の処理（必要に応じて）
    logger.info("Document Processor Service 終了")

# FastAPI アプリ初期化（ライフサイクル管理を追加）
app = FastAPI(
    title="RAG Document Processor (Ollama専用)", 
    version="2.0.0",
    lifespan=lifespan
)

# サポートするファイルタイプ
SUPPORTED_EXTENSIONS = {
    '.pdf': 'application/pdf',
    '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    '.txt': 'text/plain'
}

@app.get("/health")
async def health_check():
    """ヘルスチェックエンドポイント - モデル準備状況も含む"""
    try:
        # データベース接続テスト
        with engine.connect() as conn:
            conn.execute(text("SELECT 1"))
        db_status = "connected"
    except Exception as e:
        logger.error(f"データベース接続エラー: {e}")
        db_status = "disconnected"
    
    # Ollama埋め込みモデルの準備状況確認
    embedding_ready = await check_ollama_model_ready(EMBEDDING_MODEL)
    
    # 全体のサービス状況判定
    service_ready = embedding_ready["ready"] and db_status == "connected"
    service_status = "ready" if service_ready else "not_ready"
    
    return {
        "status": service_status,
        "service": "document-processor",
        "mode": "ollama",
        "embedding_model": {
            "name": EMBEDDING_MODEL,
            "ready": embedding_ready["ready"],
            "status": embedding_ready["status"],
            "message": embedding_ready["message"]
        },
        "database": db_status,
        "all_ready": service_ready,
        "ollama_url": OLLAMA_URL,
        "timestamp": datetime.now().isoformat()
    }

@app.get("/")
async def root():
    """ルートエンドポイント"""
    return {
        "message": "RAG Document Processor Service",
        "mode": "Ollama専用",
        "version": "2.0.0",
        "docs": "/docs"
    }

@app.get("/test-ollama")
async def test_ollama():
    """Ollama接続テスト"""
    try:
        async with httpx.AsyncClient() as client:
            # モデル一覧を取得
            response = await client.get(f"{OLLAMA_URL}/api/tags", timeout=10.0)
            if response.status_code == 200:
                models = response.json()
                
                # 埋め込みテスト
                embed_response = await client.post(
                    f"{OLLAMA_URL}/api/embeddings",
                    json={"model": "nomic-embed-text", "prompt": "test"},
                    timeout=30.0
                )
                
                return {
                    "status": "success",
                    "ollama_url": OLLAMA_URL,
                    "models_available": len(models.get("models", [])),
                    "embedding_test_status": embed_response.status_code,
                    "embedding_test_ok": embed_response.status_code == 200,
                    "models": [m["name"] for m in models.get("models", [])]
                }
            else:
                return {
                    "status": "error",
                    "ollama_url": OLLAMA_URL,
                    "error": f"Ollama API応答エラー: {response.status_code}"
                }
    except Exception as e:
        return {
            "status": "error",
            "ollama_url": OLLAMA_URL,
            "error": str(e)
        }

@app.get("/documents")
async def get_documents():
    """アップロード済みドキュメント一覧を取得"""
    try:
        with engine.connect() as conn:
            result = conn.execute(text("""
                SELECT id, filename, file_size, file_type, upload_time, 
                       processed_time, status, total_chunks, content_preview
                FROM documents 
                ORDER BY upload_time DESC
            """))
            documents = []
            for row in result:
                documents.append({
                    "id": row[0],
                    "filename": row[1],
                    "file_size": row[2],
                    "file_type": row[3],
                    "upload_time": row[4].isoformat() if row[4] else None,
                    "processed_time": row[5].isoformat() if row[5] else None,
                    "status": row[6],
                    "total_chunks": row[7],
                    "content_preview": row[8][:200] + "..." if row[8] and len(row[8]) > 200 else row[8]
                })
            return {"documents": documents, "count": len(documents)}
    except Exception as e:
        logger.error(f"ドキュメント一覧取得エラー: {e}")
        raise HTTPException(status_code=500, detail=f"ドキュメント一覧取得に失敗しました: {str(e)}")

@app.get("/documents/{document_id}")
async def get_document(document_id: int):
    """特定のドキュメント詳細を取得"""
    try:
        with engine.connect() as conn:
            result = conn.execute(text("""
                SELECT id, filename, file_path, file_size, file_type, upload_time, 
                       processed_time, status, metadata, content_preview, total_chunks
                FROM documents 
                WHERE id = :doc_id
            """), {"doc_id": document_id})
            row = result.fetchone()
            if not row:
                raise HTTPException(status_code=404, detail="ドキュメントが見つかりません")
            
            return {
                "id": row[0],
                "filename": row[1],
                "file_path": row[2],
                "file_size": row[3],
                "file_type": row[4],
                "upload_time": row[5].isoformat() if row[5] else None,
                "processed_time": row[6].isoformat() if row[6] else None,
                "status": row[7],
                "metadata": row[8],
                "content_preview": row[9],
                "total_chunks": row[10]
            }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"ドキュメント取得エラー: {e}")
        raise HTTPException(status_code=500, detail=f"ドキュメント取得に失敗しました: {str(e)}")

@app.delete("/documents/{document_id}")
async def delete_document(document_id: int):
    """ドキュメントとその関連データを削除"""
    session = SessionLocal()
    try:
        # ドキュメント情報を取得
        result = session.execute(text("SELECT filename, file_path FROM documents WHERE id = :document_id"), {"document_id": document_id})
        document = result.fetchone()
        
        if not document:
            raise HTTPException(status_code=404, detail="ドキュメントが見つかりません")
        
        filename, file_path = document
        
        # 関連チャンクを削除
        delete_chunks_result = session.execute(text("DELETE FROM document_chunks WHERE document_id = :document_id"), {"document_id": document_id})
        chunks_deleted = delete_chunks_result.rowcount
        
        # ドキュメントを削除
        delete_doc_result = session.execute(text("DELETE FROM documents WHERE id = :document_id"), {"document_id": document_id})
        
        if delete_doc_result.rowcount == 0:
            raise HTTPException(status_code=404, detail="ドキュメントが見つかりません")
        
        session.commit()
        
        # ファイルシステムからファイルを削除（オプション）
        try:
            if file_path and os.path.exists(file_path):
                os.unlink(file_path)
                logger.info(f"ファイルシステムからファイルを削除: {file_path}")
        except Exception as file_error:
            logger.warning(f"ファイル削除エラー: {file_error}")
            # ファイル削除に失敗してもDBから削除は成功とする
        
        logger.info(f"ドキュメント削除完了: {filename} (チャンク数: {chunks_deleted})")
        
        return {
            "status": "success",
            "message": f"ドキュメント '{filename}' を削除しました",
            "document_id": document_id,
            "chunks_deleted": chunks_deleted
        }
        
    except HTTPException:
        raise
    except Exception as e:
        session.rollback()
        logger.error(f"ドキュメント削除エラー: {e}")
        raise HTTPException(status_code=500, detail=f"ドキュメント削除に失敗しました: {str(e)}")
    finally:
        session.close()

def preprocess_image_for_ocr(image_path: str) -> str:
    """OCR用画像前処理 - 品質向上とノイズ除去"""
    try:
        import cv2
        import numpy as np
        from PIL import ImageEnhance, ImageFilter

        # PILで画像を開く
        img = Image.open(image_path)
        
        # グレースケール変換
        if img.mode != 'L':
            img = img.convert('L')
        
        # 解像度向上（2倍にリサイズ）
        width, height = img.size
        img = img.resize((width * 2, height * 2), Image.LANCZOS)
        
        # コントラスト強化
        enhancer = ImageEnhance.Contrast(img)
        img = enhancer.enhance(2.0)
        
        # シャープネス強化
        enhancer = ImageEnhance.Sharpness(img)
        img = enhancer.enhance(2.0)
        
        # ノイズ除去フィルタ
        img = img.filter(ImageFilter.MedianFilter(size=3))
        
        # OpenCVで更なる前処理
        img_array = np.array(img)
        
        # ガウシアンブラーでノイズ除去
        img_array = cv2.GaussianBlur(img_array, (1, 1), 0)
        
        # 二値化（OTSU手法）
        _, img_array = cv2.threshold(img_array, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        
        # モルフォロジー処理でノイズ除去
        kernel = np.ones((2, 2), np.uint8)
        img_array = cv2.morphologyEx(img_array, cv2.MORPH_CLOSE, kernel)
        
        # 前処理済み画像を一時保存
        processed_img = Image.fromarray(img_array)
        processed_path = image_path.replace('.png', '_processed.png').replace('.jpg', '_processed.jpg')
        processed_img.save(processed_path)
        
        return processed_path
        
    except Exception as e:
        logger.warning(f"画像前処理エラー: {e}")
        return image_path  # 前処理失敗時は元画像を返す

def extract_text_with_ocr(image_path: str) -> str:
    """OCRで画像からテキストを抽出（日本語対応・強化版）"""
    try:
        # 画像前処理を実行
        processed_path = preprocess_image_for_ocr(image_path)
        
        # EasyOCRを使用（前処理済み画像で）
        results = ocr_reader.readtext(processed_path)
        text = ""
        for (bbox, text_content, confidence) in results:
            if confidence > 0.3:  # 信頼度閾値を下げて、より多くのテキストを取得
                text += text_content + " "
        
        # EasyOCRで不十分な場合はTesseractで詳細処理
        if len(text.strip()) < 15:
            logger.info("EasyOCRの結果が不十分、強化版Tesseractで処理")
            img = Image.open(processed_path)
            
            # Tesseractの詳細設定（日本語最適化）
            custom_config = r'--oem 3 --psm 6 -c tessedit_char_whitelist=0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyzあいうえおかきくけこさしすせそたちつてとなにぬねのはひふへほまみむめもやゆよらりるれろわをんがぎぐげござじずぜぞだぢづでどばびぶべぼぱぴぷぺぽァィゥェォャュョッアイウエオカキクケコサシスセソタチツテトナニヌネノハヒフヘホマミムメモヤユヨラリルレロワヲンガギグゲゴザジズゼゾダヂヅデドバビブベボパピプペポ一二三四五六七八九十百千万億兆　、。（）「」【】～・'
            
            # 複数のPSMモードで試行
            for psm in [6, 7, 8, 13]:
                try:
                    config = f'--oem 3 --psm {psm}'
                    text_tesseract = pytesseract.image_to_string(img, lang='jpn+eng', config=config)
                    if len(text_tesseract.strip()) > len(text.strip()):
                        text = text_tesseract
                        logger.info(f"Tesseract PSM {psm}で最良結果: {len(text)} 文字")
                        break
                except Exception as psm_error:
                    logger.warning(f"Tesseract PSM {psm} エラー: {psm_error}")
                    continue
        
        # 一時ファイルのクリーンアップ
        if processed_path != image_path and os.path.exists(processed_path):
            try:
                os.unlink(processed_path)
            except:
                pass
        
        result_text = text.strip()
        logger.info(f"OCR結果: {len(result_text)} 文字抽出")
        return result_text
        
    except Exception as e:
        logger.error(f"OCRエラー: {e}")
        return ""

def extract_text_from_pdf(file_path: Path) -> str:
    """PDFからテキストを抽出（OCRフォールバック付き）"""
    try:
        # まず通常のテキスト抽出を試行
        with open(file_path, 'rb') as file:
            reader = pypdf.PdfReader(file)
            text = ""
            for page in reader.pages:
                page_text = page.extract_text()
                text += page_text + "\n"
        
        # テキストが十分抽出できた場合はそのまま返す
        if len(text.strip()) > 100:
            return text.strip()
        
        # テキストが少ない場合はOCRを使用
        logger.info("PDF内のテキストが少ないため、高品質OCRを使用します")
        try:
            # PDFを高解像度画像に変換してOCR処理
            images = convert_from_path(file_path, dpi=300, fmt='png')  # 高解像度でPNG変換
            ocr_text = ""
            
            logger.info(f"PDF {len(images)}ページを高品質OCRで処理中...")
            
            for i, image in enumerate(images):
                # 一時的に画像を保存（高品質PNG）
                temp_image_path = f"/tmp/pdf_page_{i+1}.png"
                image.save(temp_image_path, 'PNG', optimize=False, quality=100)
                
                # 強化されたOCRでテキスト抽出
                page_ocr_text = extract_text_with_ocr(temp_image_path)
                
                if page_ocr_text.strip():
                    ocr_text += f"ページ {i+1}:\n{page_ocr_text}\n\n"
                    logger.info(f"ページ {i+1}: {len(page_ocr_text)} 文字抽出成功")
                else:
                    logger.warning(f"ページ {i+1}: テキスト抽出失敗")
                
                # 一時ファイルを削除
                try:
                    os.unlink(temp_image_path)
                except:
                    pass
            
            if len(ocr_text.strip()) > len(text.strip()):
                logger.info(f"OCR成功: {len(ocr_text)} 文字抽出 (元: {len(text)} 文字)")
                return ocr_text.strip()
            else:
                logger.warning("OCRでも改善なし、元のテキストを返します")
                return text.strip()
        
        except Exception as ocr_error:
            logger.error(f"PDF OCRエラー: {ocr_error}")
            return text.strip()
            
    except Exception as e:
        logger.error(f"PDF読み込みエラー: {e}")
        return ""

def extract_text_from_docx(file_path: Path) -> str:
    """DOCXからテキストを抽出（LibreOfficeフォールバック）"""
    try:
        # まずpython-docxで直接テキスト抽出を試行
        doc = Document(file_path)
        direct_text = ""
        for paragraph in doc.paragraphs:
            direct_text += paragraph.text + "\n"
        
        # テーブルのテキストも抽出
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    direct_text += " | ".join(row_text) + "\n"
        
        # 直接抽出で十分なテキストが取得できた場合
        if len(direct_text.strip()) > 50:
            logger.info(f"DOCX直接抽出成功: {len(direct_text)} 文字")
            return direct_text.strip()
        
        # LibreOfficeで変換を試行
        logger.info("DOCX内のテキストが少ないため、LibreOfficeで変換を試行します")
        return extract_text_with_libreoffice(file_path, direct_text)
        
    except Exception as e:
        logger.error(f"DOCX読み込みエラー: {e}")
        # エラーの場合もLibreOfficeでフォールバック
        return extract_text_with_libreoffice(file_path, "")

def extract_text_from_pptx(file_path: Path) -> str:
    """PPTXからテキストを抽出（LibreOffice優先、OCRフォールバック）"""
    try:
        # まずpython-pptxで直接テキスト抽出を試行
        prs = Presentation(file_path)
        direct_text = ""
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = f"スライド {slide_num}:\n"
            slide_content = ""
            
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    slide_content += shape.text.strip() + "\n"
                # テーブルの処理
                elif shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            slide_content += " | ".join(row_text) + "\n"
            
            direct_text += slide_text + slide_content + "\n"
        
        # 直接抽出で十分なテキストが取得できた場合
        if len(direct_text.strip()) > 100:
            logger.info(f"PPTX直接抽出成功: {len(direct_text)} 文字")
            return direct_text.strip()
        
        # LibreOfficeでPDFに変換してからテキスト抽出
        logger.info("PPTX内のテキストが少ないため、LibreOfficeでPDF変換を試行します")
        return extract_text_with_libreoffice(file_path, direct_text)
        
    except Exception as e:
        logger.error(f"PPTX読み込みエラー: {e}")
        return ""

def extract_text_with_libreoffice(file_path: Path, fallback_text: str = "") -> str:
    """LibreOfficeを使用してファイルをPDFに変換し、テキストを抽出"""
    try:
        import subprocess
        import tempfile

        # サポートされているファイル形式をチェック
        supported_extensions = {'.pptx', '.ppt', '.docx', '.doc', '.xlsx', '.xls', '.odp', '.ods', '.odt'}
        if file_path.suffix.lower() not in supported_extensions:
            logger.warning(f"LibreOfficeでサポートされていない形式: {file_path.suffix}")
            return fallback_text
        
        # 一時ディレクトリを作成
        with tempfile.TemporaryDirectory() as temp_dir:
            output_dir = Path(temp_dir)
            
            # LibreOfficeでPDFに変換
            logger.info(f"LibreOfficeで変換中: {file_path.name} ({file_path.suffix})")
            result = subprocess.run([
                'libreoffice', '--headless', '--convert-to', 'pdf',
                '--outdir', str(output_dir), str(file_path)
            ], capture_output=True, text=True, timeout=120)
            
            if result.returncode != 0:
                logger.warning(f"LibreOffice変換失敗 (終了コード: {result.returncode})")
                if result.stderr:
                    logger.warning(f"エラー詳細: {result.stderr}")
                return fallback_text
            
            # 変換されたPDFファイルを探す
            pdf_files = list(output_dir.glob("*.pdf"))
            if not pdf_files:
                logger.warning("変換されたPDFファイルが見つかりません")
                return fallback_text
            
            pdf_path = pdf_files[0]
            logger.info(f"PDF変換成功: {pdf_path.name} (サイズ: {pdf_path.stat().st_size} bytes)")
            
            # PDFからテキスト抽出
            pdf_text = extract_text_from_pdf(pdf_path)
            
            # LibreOffice抽出結果をチェック
            if len(pdf_text.strip()) > max(len(fallback_text.strip()), 50):
                logger.info(f"LibreOffice抽出成功: {len(pdf_text)} 文字 (元: {len(fallback_text)} 文字)")
                return pdf_text
            else:
                # LibreOfficeでも改善がない場合、OCRを試行
                logger.info("LibreOfficeでもテキストが少ないため、OCRを試行します")
                return extract_text_with_ocr_from_pdf(pdf_path, fallback_text)
    
    except subprocess.TimeoutExpired:
        logger.error("LibreOffice変換がタイムアウトしました（120秒）")
        return fallback_text
    except Exception as e:
        logger.error(f"LibreOffice処理エラー: {e}")
        return fallback_text

def extract_text_with_ocr_from_pdf(pdf_path: Path, fallback_text: str = "") -> str:
    """PDFから画像を抽出してOCR処理を実行"""
    try:
        # PDFを画像に変換してOCR処理
        images = convert_from_path(pdf_path)
        ocr_text = ""
        
        for i, image in enumerate(images):
            # 一時的に画像を保存
            temp_image_path = f"/tmp/ocr_page_{i+1}.png"
            image.save(temp_image_path, 'PNG')
            
            # OCRでテキスト抽出
            page_ocr_text = extract_text_with_ocr(temp_image_path)
            if page_ocr_text.strip():
                ocr_text += f"ページ {i+1} (OCR):\n{page_ocr_text}\n\n"
            
            # 一時ファイルを削除
            os.unlink(temp_image_path)
        
        # OCR結果をチェック
        if len(ocr_text.strip()) > len(fallback_text.strip()):
            logger.info(f"OCR抽出成功: {len(ocr_text)} 文字 (元: {len(fallback_text)} 文字)")
            return ocr_text.strip()
        else:
            logger.info("OCRでも改善なし、元のテキストを返します")
            return fallback_text
    
    except Exception as e:
        logger.error(f"OCR処理エラー: {e}")
        return fallback_text

def extract_text_from_xlsx(file_path: Path) -> str:
    """XLSXからテキストを抽出（LibreOfficeフォールバック）"""
    try:
        # まずopenpyxlで直接テキスト抽出を試行
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        direct_text = ""
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            direct_text += f"シート: {sheet_name}\n"
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join([str(cell) for cell in row if cell is not None])
                if row_text.strip():
                    direct_text += row_text + "\n"
        
        # 直接抽出で十分なテキストが取得できた場合
        if len(direct_text.strip()) > 50:
            logger.info(f"XLSX直接抽出成功: {len(direct_text)} 文字")
            return direct_text.strip()
        
        # LibreOfficeで変換を試行
        logger.info("XLSX内のテキストが少ないため、LibreOfficeで変換を試行します")
        return extract_text_with_libreoffice(file_path, direct_text)
        
    except Exception as e:
        logger.error(f"XLSX読み込みエラー: {e}")
        # エラーの場合もLibreOfficeでフォールバック
        return extract_text_with_libreoffice(file_path, "")

def extract_text_from_file(file_path: Path) -> str:
    """ファイルタイプに応じてテキストを抽出"""
    extension = file_path.suffix.lower()
    
    if extension == '.pdf':
        return extract_text_from_pdf(file_path)
    elif extension == '.docx':
        return extract_text_from_docx(file_path)
    elif extension == '.pptx':
        return extract_text_from_pptx(file_path)
    elif extension == '.xlsx':
        return extract_text_from_xlsx(file_path)
    elif extension == '.txt':
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            with open(file_path, 'r', encoding='shift_jis') as f:
                return f.read()
    else:
        raise ValueError(f"サポートされていないファイル形式: {extension}")

async def check_ollama_model_ready(model_name: str) -> Dict[str, Any]:
    """Ollamaモデルの準備状況を詳細チェック"""
    try:
        async with httpx.AsyncClient() as client:
            # モデル一覧取得
            response = await client.get(f"{OLLAMA_URL}/api/tags", timeout=30.0)
            if response.status_code != 200:
                return {"ready": False, "status": "ollama_unavailable", "message": "Ollama接続不可"}
            
            models = response.json()
            available_models = [m["name"] for m in models.get("models", [])]
            
            # モデル名のマッチング（:latest タグも考慮）
            model_found = False
            matched_model = model_name
            
            for available_model in available_models:
                # 完全一致
                if model_name == available_model:
                    model_found = True
                    matched_model = available_model
                    break
                # :latest タグ付きでマッチ
                elif f"{model_name}:latest" == available_model:
                    model_found = True
                    matched_model = available_model
                    break
                # ベース名でマッチ（逆パターン）
                elif available_model.replace(":latest", "") == model_name:
                    model_found = True
                    matched_model = available_model
                    break
            
            if not model_found:
                return {"ready": False, "status": "model_not_found", "message": f"モデル '{model_name}' が見つかりません"}
            
            # 埋め込みモデルの場合はembeddingsエンドポイントでテスト
            if "embed" in model_name.lower():
                test_response = await client.post(
                    f"{OLLAMA_URL}/api/embeddings",
                    json={"model": matched_model, "prompt": "test"},
                    timeout=60.0
                )
            else:
                # LLMモデルの場合はgenerateエンドポイントでテスト
                test_response = await client.post(
                    f"{OLLAMA_URL}/api/generate",
                    json={"model": matched_model, "prompt": "test", "stream": False},
                    timeout=60.0
                )
            
            if test_response.status_code == 200:
                return {"ready": True, "status": "ready", "message": f"モデル '{matched_model}' 準備完了"}
            else:
                return {"ready": False, "status": "model_loading", "message": f"モデル '{matched_model}' 読み込み中"}
                
    except Exception as e:
        logger.error(f"モデル準備状況確認エラー: {e}")
        return {"ready": False, "status": "error", "message": f"エラー: {str(e)}"}

async def get_embeddings_ollama(text: str) -> List[float]:
    """Ollamaで埋め込みベクトルを生成"""
    try:
        async with httpx.AsyncClient() as client:
            response = await client.post(
                f"{OLLAMA_URL}/api/embeddings",
                json={"model": EMBEDDING_MODEL, "prompt": text},
                timeout=30.0
            )
            if response.status_code == 200:
                result = response.json()
                return result.get("embedding", [])
            else:
                logger.error(f"Ollama埋め込み生成エラー: {response.status_code}")
                # フォールバック: ダミーベクトル（768次元）
                return [0.0] * 768
    except Exception as e:
        logger.error(f"Ollama接続エラー: {e}")
        # フォールバック: ダミーベクトル（768次元）
        return [0.0] * 768

@app.post("/upload")
async def upload_document(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    """ドキュメントアップロード"""
    try:
        # ファイル検証
        if not file.filename:
            raise HTTPException(status_code=400, detail="ファイル名が必要です")
        
        file_extension = Path(file.filename).suffix.lower()
        if file_extension not in SUPPORTED_EXTENSIONS:
            raise HTTPException(
                status_code=400, 
                detail=f"サポートされていないファイル形式: {file_extension}"
            )
        
        # ファイル保存
        UPLOAD_DIR.mkdir(exist_ok=True)
        file_path = UPLOAD_DIR / file.filename
        
        with open(file_path, "wb") as buffer:
            content = await file.read()
            buffer.write(content)
        
        # バックグラウンドでドキュメント処理
        background_tasks.add_task(process_document, file_path, file.filename)
        
        return {
            "status": "success",
            "message": "ファイルアップロード完了。バックグラウンドで処理中...",
            "filename": file.filename,
            "file_size": len(content),
            "mode": "Ollama"
        }
        
    except Exception as e:
        logger.error(f"アップロードエラー: {e}")
        raise HTTPException(status_code=500, detail=str(e))

async def process_document(file_path: Path, filename: str):
    """ドキュメント処理（非同期）"""
    try:
        logger.info(f"ドキュメント処理開始: {filename}")
        
        # テキスト抽出
        content = extract_text_from_file(file_path)
        if not content.strip():
            logger.warning(f"テキストが抽出できませんでした: {filename}")
            return
        
        # チャンク分割（改良版）- より小さなチャンクで精度向上
        chunks = [content[i:i+600] for i in range(0, len(content), 400)]
        
        # データベース保存準備
        session = SessionLocal()
        try:
            # ドキュメント情報保存
            doc_id = await save_document_info(session, filename, file_path, content)
            
            # チャンク処理
            for i, chunk in enumerate(chunks):
                embedding = await get_embeddings_ollama(chunk)
                await save_document_chunk(session, doc_id, i, chunk, embedding)
            
            # ドキュメントステータスを完了に更新
            session.execute(text("""
                UPDATE documents 
                SET status = 'processed', processed_time = CURRENT_TIMESTAMP, total_chunks = :total_chunks
                WHERE id = :doc_id
            """), {"doc_id": doc_id, "total_chunks": len(chunks)})
            
            session.commit()
            logger.info(f"ドキュメント処理完了: {filename} ({len(chunks)}チャンク)")
            
        except Exception as e:
            session.rollback()
            logger.error(f"データベース保存エラー: {e}")
            # エラー時のステータス更新
            try:
                if 'doc_id' in locals():
                    session.execute(text("""
                        UPDATE documents 
                        SET status = 'error', processed_time = CURRENT_TIMESTAMP
                        WHERE id = :doc_id
                    """), {"doc_id": doc_id})
                    session.commit()
            except Exception as update_error:
                logger.error(f"エラーステータス更新失敗: {update_error}")
        finally:
            session.close()
            
    except Exception as e:
        logger.error(f"ドキュメント処理エラー: {e}")

async def save_document_info(session, filename: str, file_path: Path, content: str) -> int:
    """ドキュメント情報をデータベースに保存"""
    file_size = file_path.stat().st_size
    file_type = SUPPORTED_EXTENSIONS.get(file_path.suffix.lower(), "unknown")
    
    # ファイル名の長さ制限（255文字）
    safe_filename = filename[:255] if len(filename) > 255 else filename
    
    # ファイルタイプの長さ制限（100文字）
    safe_file_type = file_type[:100] if len(file_type) > 100 else file_type
    
    # コンテンツプレビューの長さ制限（500文字）
    safe_content_preview = content[:500] if len(content) > 500 else content
    
    try:
        result = session.execute(text("""
            INSERT INTO documents (filename, file_path, file_size, file_type, content_preview, status)
            VALUES (:filename, :file_path, :file_size, :file_type, :content_preview, 'processing')
            RETURNING id
        """), {
            "filename": safe_filename,
            "file_path": str(file_path),
            "file_size": file_size,
            "file_type": safe_file_type,
            "content_preview": safe_content_preview
        })
        
        doc_id = result.scalar()
        logger.info(f"ドキュメント情報保存完了: {safe_filename} (ID: {doc_id})")
        return doc_id
        
    except Exception as e:
        logger.error(f"ドキュメント情報保存エラー: {e}")
        logger.error(f"パラメータ: filename={safe_filename}, file_type={safe_file_type}")
        raise

async def save_document_chunk(session, doc_id: int, chunk_index: int, content: str, embedding: List[float]):
    """ドキュメントチャンクをデータベースに保存"""
    session.execute(text("""
        INSERT INTO document_chunks (document_id, chunk_index, content, embedding)
        VALUES (:document_id, :chunk_index, :content, :embedding)
        ON CONFLICT (document_id, chunk_index) 
        DO UPDATE SET content = EXCLUDED.content, embedding = EXCLUDED.embedding
    """), {
        "document_id": doc_id,
        "chunk_index": chunk_index,
        "content": content,
        "embedding": embedding
    })

if __name__ == "__main__":
    uvicorn.run(
        "app:app", 
        host="0.0.0.0", 
        port=8001,
        reload=False,
        log_level="info"
    )
